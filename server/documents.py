"""Document generation: Word (.docx) and PowerPoint (.pptx) from agent outputs."""

from __future__ import annotations

import json
import logging
import time
import uuid
from pathlib import Path
from typing import Any

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

DOWNLOAD_DIR = Path(__file__).resolve().parent.parent / "data" / "downloads"
MAX_FILES = 200  # max stored files before cleanup

_files: dict[str, tuple[Path, float]] = {}  # id -> (path, timestamp)


def store_file(path: Path) -> str:
    file_id = uuid.uuid4().hex[:12]
    _files[file_id] = (path, time.time())
    _cleanup_old_files()
    return file_id


def get_file(file_id: str) -> Path | None:
    entry = _files.get(file_id)
    if entry:
        return entry[0]
    # fallback: check disk
    if DOWNLOAD_DIR.exists():
        for f in DOWNLOAD_DIR.iterdir():
            if file_id in f.name:
                return f
    return None


def _cleanup_old_files() -> None:
    """Remove oldest files when storage exceeds MAX_FILES."""
    if len(_files) <= MAX_FILES:
        return
    sorted_entries = sorted(_files.items(), key=lambda x: x[1][1])
    to_remove = sorted_entries[:len(_files) - MAX_FILES]
    for fid, (path, _) in to_remove:
        try:
            if path.exists():
                path.unlink()
        except OSError:
            pass
        del _files[fid]
    logger.info("Cleaned up %d old download files", len(to_remove))


def cleanup_downloads_on_startup() -> None:
    """Remove download files older than 24 hours at startup."""
    if not DOWNLOAD_DIR.exists():
        return
    cutoff = time.time() - 86400
    removed = 0
    for f in DOWNLOAD_DIR.iterdir():
        if f.stat().st_mtime < cutoff:
            try:
                f.unlink()
                removed += 1
            except OSError:
                pass
    if removed:
        logger.info("Startup cleanup: removed %d old download files", removed)


# ── Word Document ────────────────────────────────────────────────────

def generate_docx(task: str, context: dict[str, Any]) -> str:
    """Generate a Word document from all agent outputs. Returns file_id."""
    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    doc = Document()

    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)

    # Title
    title = doc.add_heading(task, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("Autonomous Startup Builder Report", style="Subtitle").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("")

    # Table of Contents placeholder
    doc.add_heading("Table of Contents", level=1)
    toc_items = [
        "1. Executive Summary",
        "2. Market Research",
        "3. Competitor Analysis",
        "4. Product Requirements",
        "5. System Architecture",
        "6. Go-to-Market Strategy",
        "7. MVP Code Overview",
        "8. Investor Pitch",
    ]
    for item in toc_items:
        doc.add_paragraph(item)
    doc.add_page_break()

    # 1. Executive Summary (from research)
    research = context.get("ResearchAgent", {})
    doc.add_heading("1. Executive Summary", level=1)
    market = research.get("market_overview", {})
    if market:
        doc.add_paragraph(f"Market Size: {market.get('market_size', 'N/A')}")
        doc.add_paragraph(f"Growth Rate: {market.get('growth_rate', 'N/A')}")
        _add_list(doc, "Key Segments", market.get("key_segments", []))
    doc.add_page_break()

    # 2. Market Research
    doc.add_heading("2. Market Research", level=1)
    audience = research.get("target_audience", {})
    if audience:
        primary = audience.get("primary", {})
        doc.add_heading("Primary Audience", level=2)
        doc.add_paragraph(f"Segment: {primary.get('segment', 'N/A')}")
        _add_list(doc, "Pain Points", primary.get("pain_points", []))
        if primary.get("buying_behavior"):
            doc.add_paragraph(f"Buying Behavior: {primary['buying_behavior']}")
        secondary = audience.get("secondary", {})
        if secondary:
            doc.add_heading("Secondary Audience", level=2)
            doc.add_paragraph(f"Segment: {secondary.get('segment', 'N/A')}")
            _add_list(doc, "Pain Points", secondary.get("pain_points", []))
    _add_list(doc, "Industry Trends", research.get("trends", []))
    _add_list(doc, "Key Risks", research.get("risks", []))
    doc.add_page_break()

    # 3. Competitor Analysis
    analysis = context.get("AnalysisAgent", {})
    doc.add_heading("3. Competitor Analysis", level=1)
    competitors = analysis.get("competitors", [])
    if competitors:
        table = doc.add_table(rows=1, cols=4)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "Competitor"
        hdr[1].text = "Strengths"
        hdr[2].text = "Weaknesses"
        hdr[3].text = "Pricing"
        for c in competitors:
            row = table.add_row().cells
            row[0].text = str(c.get("name", ""))
            row[1].text = ", ".join(str(x) for x in c.get("strengths", []))
            row[2].text = ", ".join(str(x) for x in c.get("weaknesses", []))
            row[3].text = str(c.get("pricing", "N/A"))

    swot = analysis.get("swot", {})
    if swot:
        doc.add_heading("SWOT Analysis", level=2)
        for key in ["strengths", "weaknesses", "opportunities", "threats"]:
            _add_list(doc, key.capitalize(), swot.get(key, []))

    pos = analysis.get("positioning", {})
    if pos:
        doc.add_heading("Positioning", level=2)
        doc.add_paragraph(f"Differentiation: {pos.get('differentiation', 'N/A')}")
        doc.add_paragraph(f"Value Proposition: {pos.get('value_proposition', 'N/A')}")
    doc.add_page_break()

    # 4. Product Requirements
    prd = context.get("ProductAgent", {})
    doc.add_heading("4. Product Requirements", level=1)
    if prd.get("title"):
        doc.add_paragraph(prd["title"], style="Intense Quote")
    if prd.get("vision"):
        doc.add_paragraph(f"Vision: {prd['vision']}")
    _add_list(doc, "Success Metrics", prd.get("success_metrics", []))

    personas = prd.get("user_personas", [])
    if personas:
        doc.add_heading("User Personas", level=2)
        for p in personas:
            doc.add_heading(f"{p.get('name', '')} — {p.get('role', '')}", level=3)
            _add_list(doc, "Goals", p.get("goals", []))
            _add_list(doc, "Frustrations", p.get("frustrations", []))

    features = prd.get("features", {})
    mvp = features.get("mvp", [])
    if mvp:
        doc.add_heading("MVP Features", level=2)
        table = doc.add_table(rows=1, cols=3)
        table.style = "Light Grid Accent 1"
        hdr = table.rows[0].cells
        hdr[0].text = "Feature"
        hdr[1].text = "Priority"
        hdr[2].text = "Description"
        for f in mvp:
            row = table.add_row().cells
            row[0].text = f.get("name", "")
            row[1].text = f.get("priority", "")
            row[2].text = f.get("description", "")

    nfr = prd.get("non_functional", {})
    if nfr:
        doc.add_heading("Non-Functional Requirements", level=2)
        for k, v in nfr.items():
            doc.add_paragraph(f"{k.replace('_', ' ').capitalize()}: {v}")

    timeline = prd.get("timeline", {})
    if timeline:
        doc.add_heading("Timeline", level=2)
        for phase, desc in timeline.items():
            doc.add_paragraph(f"{phase.replace('_', ' ').capitalize()}: {desc}")
    doc.add_page_break()

    # 5. System Architecture
    arch = context.get("ArchitectAgent", {})
    doc.add_heading("5. System Architecture", level=1)
    if arch.get("architecture_style"):
        doc.add_paragraph(f"Style: {arch['architecture_style']}")
    if arch.get("system_diagram"):
        doc.add_paragraph(f"System Diagram: {arch['system_diagram']}")

    tech = arch.get("tech_stack", {})
    if tech:
        doc.add_heading("Tech Stack", level=2)
        for layer, items in tech.items():
            doc.add_heading(layer.replace("_", " ").capitalize(), level=3)
            if isinstance(items, dict):
                for k, v in items.items():
                    doc.add_paragraph(f"{k}: {v}")
            else:
                doc.add_paragraph(str(items))

    services = arch.get("services", [])
    if services:
        doc.add_heading("Services", level=2)
        for s in services:
            doc.add_heading(s.get("name", ""), level=3)
            doc.add_paragraph(f"Responsibility: {s.get('responsibility', '')}")
            doc.add_paragraph(f"Technology: {s.get('tech', '')}")
            _add_list(doc, "Endpoints", s.get("endpoints", []))
    doc.add_page_break()

    # 6. Go-to-Market
    marketing = context.get("MarketingAgent", {})
    doc.add_heading("6. Go-to-Market Strategy", level=1)
    messaging = marketing.get("messaging", {})
    if messaging:
        doc.add_heading("Messaging", level=2)
        if messaging.get("tagline"):
            doc.add_paragraph(messaging["tagline"], style="Intense Quote")
        if messaging.get("elevator_pitch"):
            doc.add_paragraph(messaging["elevator_pitch"])
        _add_list(doc, "Key Messages", messaging.get("key_messages", []))

    launch = marketing.get("launch_plan", [])
    if launch:
        doc.add_heading("Launch Plan", level=2)
        for phase in launch:
            doc.add_heading(f"{phase.get('phase', '')} ({phase.get('timeline', '')})", level=3)
            _add_list(doc, "Activities", phase.get("activities", []))
            _add_list(doc, "KPIs", phase.get("kpis", []))

    pricing = marketing.get("pricing", {})
    if pricing:
        doc.add_heading("Pricing", level=2)
        doc.add_paragraph(f"Model: {pricing.get('model', '')}")
        tiers = pricing.get("tiers", [])
        if tiers:
            table = doc.add_table(rows=1, cols=3)
            table.style = "Light Grid Accent 1"
            hdr = table.rows[0].cells
            hdr[0].text = "Tier"
            hdr[1].text = "Price"
            hdr[2].text = "Features"
            for t in tiers:
                row = table.add_row().cells
                row[0].text = t.get("name", "")
                row[1].text = t.get("price", "")
                row[2].text = t.get("features", "")
    doc.add_page_break()

    # 7. MVP Code Overview
    code = context.get("CodeAgent", {})
    doc.add_heading("7. MVP Code Overview", level=1)
    structure = code.get("project_structure", [])
    if isinstance(structure, list):
        _add_list(doc, "Project Structure", structure)
    _add_list(doc, "Setup Instructions", code.get("setup_instructions", []))

    for section_name, section_key in [("Backend", "backend"), ("Frontend", "frontend"), ("Deployment", "deployment")]:
        section = code.get(section_key, {})
        if section and isinstance(section, dict):
            doc.add_heading(section_name, level=2)
            for filename, content in section.items():
                doc.add_heading(filename, level=3)
                if isinstance(content, str) and len(content) > 20:
                    p = doc.add_paragraph()
                    run = p.add_run(content[:2000])
                    run.font.name = "Consolas"
                    run.font.size = Pt(8)
                    if len(content) > 2000:
                        doc.add_paragraph("... (truncated)")
    doc.add_page_break()

    # 8. Investor Pitch
    pitch = context.get("PitchAgent", {})
    if "raw_output" in pitch and "slides" not in pitch:
        raw = pitch["raw_output"]
        start = raw.find("{")
        end = raw.rfind("}")
        if start != -1 and end > start:
            try:
                pitch = json.loads(raw[start:end + 1])
            except (json.JSONDecodeError, Exception):
                pass
    doc.add_heading("8. Investor Pitch", level=1)
    slides = pitch.get("slides", [])
    for slide in slides:
        title_text = slide.get("title", f"Slide {slide.get('number', '')}")
        doc.add_heading(title_text, level=2)
        content = slide.get("content", {})
        if isinstance(content, dict):
            for k, v in content.items():
                if isinstance(v, list):
                    _add_list(doc, k.replace("_", " ").capitalize(), v)
                elif isinstance(v, dict):
                    for kk, vv in v.items():
                        doc.add_paragraph(f"{kk}: {vv}")
                else:
                    doc.add_paragraph(f"{k.replace('_', ' ').capitalize()}: {v}")

    faq = pitch.get("investor_faq", [])
    if faq:
        doc.add_heading("Investor FAQ", level=2)
        for item in faq:
            doc.add_paragraph(f"Q: {item.get('question', '')}", style="List Bullet")
            doc.add_paragraph(f"A: {item.get('answer', '')}")

    projections = pitch.get("financial_projections", {})
    if projections:
        doc.add_heading("Financial Projections", level=2)
        for year, data in projections.items():
            doc.add_heading(year.replace("_", " ").capitalize(), level=3)
            if isinstance(data, dict):
                for k, v in data.items():
                    doc.add_paragraph(f"{k.capitalize()}: {v}")

    # Save
    filename = f"startup_plan_{uuid.uuid4().hex[:8]}.docx"
    path = DOWNLOAD_DIR / filename
    doc.save(str(path))
    return store_file(path)


def _add_list(doc: Document, heading: str, items: list) -> None:
    if not items:
        return
    doc.add_paragraph(heading + ":", style="List Bullet")
    for item in items:
        doc.add_paragraph(str(item), style="List Bullet 2")


# ── PowerPoint ───────────────────────────────────────────────────────

def generate_pptx(task: str, context: dict[str, Any]) -> str:
    """Generate a PowerPoint from PitchAgent output. Returns file_id."""
    DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    pitch = context.get("PitchAgent", {})

    # If pitch data came as raw_output (JSON parse failed), try to re-parse
    if "raw_output" in pitch and "slides" not in pitch:
        import re as _re
        raw = pitch["raw_output"]
        start = raw.find("{")
        end = raw.rfind("}")
        if start != -1 and end > start:
            try:
                pitch = json.loads(raw[start:end + 1])
            except (json.JSONDecodeError, Exception):
                pass

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, PptxRGB(0x09, 0x09, 0x0B))

    title_box = slide.shapes.add_textbox(PptxInches(1.5), PptxInches(2), PptxInches(10), PptxInches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = task
    p.font.size = PptxPt(44)
    p.font.color.rgb = PptxRGB(0xFA, 0xFA, 0xFA)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    subtitle = tf.add_paragraph()
    subtitle.text = "Startup Blueprint"
    subtitle.font.size = PptxPt(20)
    subtitle.font.color.rgb = PptxRGB(0xA1, 0xA1, 0xAA)
    subtitle.alignment = PP_ALIGN.CENTER

    # Content slides from PitchAgent
    slides_data = pitch.get("slides", [])
    for s in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, PptxRGB(0x09, 0x09, 0x0B))

        # Title
        title_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(0.5), PptxInches(11), PptxInches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = s.get("title", "")
        p.font.size = PptxPt(36)
        p.font.color.rgb = PptxRGB(0xFA, 0xFA, 0xFA)
        p.font.bold = True

        # Content
        content = s.get("content", {})
        content_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.8), PptxInches(11), PptxInches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        first = True
        for key, val in content.items():
            if isinstance(val, list):
                if not first:
                    spacer = tf.add_paragraph()
                    spacer.space_before = PptxPt(12)
                for item in val:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = f"  {item}"
                    p.font.size = PptxPt(18)
                    p.font.color.rgb = PptxRGB(0xE4, 0xE4, 0xE7)
                    p.space_before = PptxPt(6)
            elif isinstance(val, dict):
                for kk, vv in val.items():
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = f"{kk}: {vv}"
                    p.font.size = PptxPt(16)
                    p.font.color.rgb = PptxRGB(0xA1, 0xA1, 0xAA)
                    p.space_before = PptxPt(4)
            else:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = str(val)
                p.font.size = PptxPt(20)
                p.font.color.rgb = PptxRGB(0xFA, 0xFA, 0xFA)
                p.space_before = PptxPt(8)

    # Financial projections slide
    projections = pitch.get("financial_projections", {})
    if projections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, PptxRGB(0x09, 0x09, 0x0B))

        title_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(0.5), PptxInches(11), PptxInches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Financial Projections"
        p.font.size = PptxPt(36)
        p.font.color.rgb = PptxRGB(0xFA, 0xFA, 0xFA)
        p.font.bold = True

        content_box = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.8), PptxInches(11), PptxInches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        first = True
        for year, data in projections.items():
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = year.replace("_", " ").upper()
            p.font.size = PptxPt(24)
            p.font.color.rgb = PptxRGB(0xFA, 0xFA, 0xFA)
            p.font.bold = True
            p.space_before = PptxPt(16)
            if isinstance(data, dict):
                for k, v in data.items():
                    p = tf.add_paragraph()
                    p.text = f"    {k.capitalize()}: {v}"
                    p.font.size = PptxPt(16)
                    p.font.color.rgb = PptxRGB(0xA1, 0xA1, 0xAA)

    # Closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PptxRGB(0x09, 0x09, 0x0B))
    title_box = slide.shapes.add_textbox(PptxInches(1.5), PptxInches(2.5), PptxInches(10), PptxInches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = PptxPt(48)
    p.font.color.rgb = PptxRGB(0xFA, 0xFA, 0xFA)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = task
    sub.font.size = PptxPt(18)
    sub.font.color.rgb = PptxRGB(0x52, 0x52, 0x5B)
    sub.alignment = PP_ALIGN.CENTER

    filename = f"startup_pitch_{uuid.uuid4().hex[:8]}.pptx"
    path = DOWNLOAD_DIR / filename
    prs.save(str(path))
    return store_file(path)


def _set_slide_bg(slide, color: PptxRGB) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color
